"""Generate PowerPoint presentations from survey data."""
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


RAMS_RED = RGBColor(0x8B, 0x1A, 0x1A)
RAMS_DARK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x27, 0xAE, 0x60)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)
RED = RGBColor(0xE7, 0x4C, 0x3C)


def _add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    for shape in slide.placeholders:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RAMS_DARK


def _add_content_slide(prs, title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RAMS_RED
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.space_after = Pt(6)


def _add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RAMS_RED

    # Table
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), Inches(1.2), Inches(9), Inches(0.4 * n_rows)
    )
    table = table_shape.table

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RAMS_RED

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)


def generate_pptx(survey_label, sections_data):
    """Generate a PowerPoint file and return as bytes.

    sections_data is a dict with optional keys:
        overview: {responses, likert_qs, yn_qs, open_qs, grades: dict}
        agreement: [(question, pct), ...]
        categories: [(category, score, pct), ...]
        benchmarks: [(indicator, rams_pct, mwahs_pct, diff), ...]
        insights: [text, ...]
        at_risk: [(label, count, total, pct), ...]
        actions: [(priority, action, finding), ...]
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    _add_title_slide(prs, "RAMS CARE Survey Report", survey_label)

    # Slide 2: Overview
    if "overview" in sections_data:
        ov = sections_data["overview"]
        bullets = [
            f"Total Responses: {ov['responses']}",
            f"Likert Questions: {ov['likert_qs']}",
            f"Yes/No Questions: {ov['yn_qs']}",
            f"Open-Ended Questions: {ov['open_qs']}",
        ]
        if ov.get("grades"):
            for grade, count in ov["grades"].items():
                bullets.append(f"{grade}: {count} responses")
        _add_content_slide(prs, "Overview", bullets)

    # Slide 3: Agreement levels (top/bottom 5)
    if "agreement" in sections_data:
        items = sections_data["agreement"]
        top5 = sorted(items, key=lambda x: x[1], reverse=True)[:5]
        bottom5 = sorted(items, key=lambda x: x[1])[:5]
        bullets = ["TOP 5 (Most Positive):"] + [f"  {q}: {p:.0f}%" for q, p in top5]
        bullets += ["", "BOTTOM 5 (Least Positive):"] + [f"  {q}: {p:.0f}%" for q, p in bottom5]
        _add_content_slide(prs, "Agreement Levels", bullets)

    # Slide 4: Categories
    if "categories" in sections_data:
        _add_table_slide(
            prs, "CARE Category Scores",
            ["Category", "Score (1-4)", "% Positive"],
            [(cat, f"{score:.2f}", f"{pct:.0f}%") for cat, score, pct in sections_data["categories"]],
        )

    # Slide 5: Benchmarks
    if "benchmarks" in sections_data and sections_data["benchmarks"]:
        _add_table_slide(
            prs, "MetroWest Benchmark Comparison",
            ["Indicator", "RAMS", "MWAHS", "Difference"],
            [(ind, f"{r:.0f}%", f"{m}%", f"{d:+.1f}%") for ind, r, m, d in sections_data["benchmarks"]],
        )

    # Slide 6: Key insights
    if "insights" in sections_data and sections_data["insights"]:
        # Strip markdown bold
        clean = [t.replace("**", "") for t in sections_data["insights"][:8]]
        _add_content_slide(prs, "Key Insights", clean)

    # Slide 7: At-risk
    if "at_risk" in sections_data and sections_data["at_risk"]:
        _add_table_slide(
            prs, "At-Risk Indicators",
            ["Indicator", "Count", "% of Respondents"],
            [(label, str(count), f"{pct:.1f}%") for label, count, total, pct in sections_data["at_risk"]],
        )

    # Slide 8: Recommended actions
    if "actions" in sections_data and sections_data["actions"]:
        bullets = []
        for priority, action, finding in sections_data["actions"][:6]:
            bullets.append(f"[{priority.upper()}] {action}")
            bullets.append(f"  → {finding.replace('**', '')}")
        _add_content_slide(prs, "Recommended Actions", bullets)

    # Save to bytes
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
